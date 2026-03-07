#!/usr/bin/env python
"""Create a PowerPoint slide for an ABB or SBB building block.

Reads components.png and summary.png from a building block folder and
produces a single-slide .pptx with the images positioned for a 16:9
widescreen presentation (13.333in x 7.500in).

Layout matches the reference template:
  - Title bar at top (full width)
  - Components diagram on the left half
  - Summary panel on the right half

Usage:
    python create-building-block-slide.py <block-folder> [--output <path>]

Examples:
    python create-building-block-slide.py building-blocks/architecture-building-blocks/AB-008/
    python create-building-block-slide.py building-blocks/solution-building-blocks/SB-011/ --output slides/SB-011.pptx
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Emu(12192000)   # 13.333in
SLIDE_HEIGHT = Emu(6858000)   # 7.500in

# Title bar position and size
TITLE_LEFT = Emu(287783)      # 0.315in
TITLE_TOP = Emu(173111)       # 0.189in
TITLE_WIDTH = Emu(10515600)   # 11.500in
TITLE_HEIGHT = Emu(558152)    # 0.610in

# Components image (left half)
COMP_LEFT = Emu(223637)       # 0.245in
COMP_TOP = Emu(720989)        # 0.788in
COMP_MAX_WIDTH = Emu(5962542) # 6.521in
COMP_MAX_HEIGHT = Emu(5491194) # 6.005in

# Summary image (right half)
SUMM_LEFT = Emu(6516759)      # 7.127in
SUMM_TOP = Emu(669618)        # 0.732in
SUMM_MAX_WIDTH = Emu(5387458) # 5.892in
SUMM_MAX_HEIGHT = Emu(5970613) # 6.530in

# Brand colours
CHARCOAL = RGBColor(0x1A, 0x1A, 0x2E)    # 1.1
OBSIDIAN = RGBColor(0x2C, 0x30, 0x38)     # 3.1


def detect_block_type(folder: Path) -> str:
    """Detect whether this is an ABB or SBB from the index.md front matter."""
    index = folder / "index.md"
    if not index.exists():
        return "Building Block"
    text = index.read_text(encoding="utf-8")
    if re.search(r"\bSBB ID\b", text):
        return "SBB"
    if re.search(r"\bABB ID\b", text):
        return "ABB"
    return "Building Block"


def extract_title(folder: Path) -> str:
    """Extract the block title from index.md front matter."""
    index = folder / "index.md"
    if not index.exists():
        return folder.name
    text = index.read_text(encoding="utf-8")
    match = re.search(r'^title:\s*["\']?(.+?)["\']?\s*$', text, re.MULTILINE)
    if match:
        return match.group(1)
    return folder.name


def fit_image(img_path: Path, max_width: int, max_height: int):
    """Calculate scaled dimensions that fit within max bounds, preserving aspect ratio."""
    from PIL import Image
    with Image.open(img_path) as img:
        img_w, img_h = img.size

    scale_w = max_width / img_w
    scale_h = max_height / img_h
    scale = min(scale_w, scale_h)

    return int(img_w * scale), int(img_h * scale)


def create_slide(block_folder: Path, output_path: Path):
    """Create the PowerPoint slide."""
    components_png = block_folder / "components.png"
    summary_png = block_folder / "summary.png"

    if not components_png.exists():
        print(f"ERROR: {components_png} not found")
        sys.exit(1)
    if not summary_png.exists():
        print(f"ERROR: {summary_png} not found")
        sys.exit(1)

    block_type = detect_block_type(block_folder)
    title_text = extract_title(block_folder)

    # Prefix with block type if not already present
    if block_type != "Building Block" and not title_text.startswith(block_type):
        title_text = f"{block_type}: {title_text}"

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title text box
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = "Helvetica"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = CHARCOAL

    # Components image (left) - fit within bounds
    try:
        cw, ch = fit_image(components_png, COMP_MAX_WIDTH, COMP_MAX_HEIGHT)
        slide.shapes.add_picture(
            str(components_png), COMP_LEFT, COMP_TOP, Emu(cw), Emu(ch)
        )
    except ImportError:
        # Pillow not available; use max dimensions directly
        slide.shapes.add_picture(
            str(components_png), COMP_LEFT, COMP_TOP, COMP_MAX_WIDTH, COMP_MAX_HEIGHT
        )

    # Summary image (right) - fit within bounds
    try:
        sw, sh = fit_image(summary_png, SUMM_MAX_WIDTH, SUMM_MAX_HEIGHT)
        slide.shapes.add_picture(
            str(summary_png), SUMM_LEFT, SUMM_TOP, Emu(sw), Emu(sh)
        )
    except ImportError:
        slide.shapes.add_picture(
            str(summary_png), SUMM_LEFT, SUMM_TOP, SUMM_MAX_WIDTH, SUMM_MAX_HEIGHT
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Created: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Create a PowerPoint slide for an ABB or SBB building block."
    )
    parser.add_argument(
        "block_folder",
        type=Path,
        help="Path to the building block folder (must contain components.png and summary.png)",
    )
    parser.add_argument(
        "--output", "-o",
        type=Path,
        default=None,
        help="Output .pptx path (default: <block-folder>/components-and-summary.pptx)",
    )
    args = parser.parse_args()

    folder = args.block_folder.resolve()
    if not folder.is_dir():
        print(f"ERROR: {folder} is not a directory")
        sys.exit(1)

    output = args.output
    if output is None:
        output = folder / "components-and-summary.pptx"
    output = output.resolve()

    create_slide(folder, output)


if __name__ == "__main__":
    main()
